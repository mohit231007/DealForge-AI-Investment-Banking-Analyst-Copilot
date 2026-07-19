import json
from pathlib import Path
from zipfile import ZipFile

from openpyxl import load_workbook
from pptx import Presentation

from dealforge_portfolio import generate_portfolio_pack


EXPECTED_ARTIFACTS = {
    "01_company_profile.md",
    "02_investment_memo.md",
    "03_source_index.json",
    "04_synthetic_document_chunks.jsonl",
    "05_due_diligence_red_flags.md",
    "06_audit_trail.md",
    "07_citation_map.json",
    "08_source_map.json",
    "09_validation_report.json",
    "10_valuation_model.xlsx",
    "11_pitchbook.pptx",
    "12_financial_extract.json",
    "13_comps_analysis.json",
    "14_precedent_transactions.json",
    "15_consistency_checks.json",
    "16_investment_committee_memo.md",
    "17_due_diligence_request_list.md",
    "18_source_confidence_ladder.json",
    "19_valuation_input_audit.json",
    "20_valuation_adjustment_recommendations.json",
    "21_corrected_peer_comps.csv",
    "22_corrected_precedent_transactions.csv",
    "23_adjusted_valuation_summary.json",
    "00_deal_pack_bundle_manifest.json",
}


def test_portfolio_pack_generates_dealforge_grade_public_artifacts(tmp_path: Path) -> None:
    result = generate_portfolio_pack("sample_data", tmp_path / "pack")
    summary = result.summary

    assert summary["company"] == "AstraFood Technologies Ltd"
    assert summary["included_peer_count"] == 4
    assert summary["excluded_peer_count"] == 1
    assert summary["included_precedent_count"] == 4
    assert summary["excluded_precedent_count"] == 1
    assert summary["adjusted_medians"]["comps_ev_revenue"] == 6.4
    assert summary["adjusted_medians"]["comps_ev_ebitda"] == 60.0
    assert summary["adjusted_medians"]["precedent_ev_revenue"] == 4.65
    assert summary["adjusted_implied_enterprise_values"]["range"] == {
        "low": 45000.0,
        "high": 80000.0,
    }
    assert summary["source_verification_status"] == "NEEDS_REVIEW"
    assert summary["unit_normalization_status"] == "NEEDS_NORMALIZATION"
    assert summary["artifact_count"] == 24

    names = {path.name for path in result.artifacts}
    assert names == EXPECTED_ARTIFACTS

    adjusted = json.loads((result.output_dir / "23_adjusted_valuation_summary.json").read_text(encoding="utf-8"))
    assert adjusted["status"] == "PASS_WITH_RECOMMENDATIONS"
    assert len(adjusted["included_peer_rows"]) == 4
    assert len(adjusted["excluded_peer_rows"]) == 1


def test_office_outputs_open_and_bundle_contains_all_artifacts(tmp_path: Path) -> None:
    result = generate_portfolio_pack("sample_data", tmp_path / "pack")

    workbook = load_workbook(result.output_dir / "10_valuation_model.xlsx", data_only=False)
    assert workbook.sheetnames == [
        "01_Read_Me",
        "02_Source_Index",
        "03_Financial_Extract",
        "04_Historical_Financials",
        "05_DCF_Assumptions",
        "06_DCF_Model",
        "07_Comps_QA",
        "08_Precedent_QA",
        "09_Source_Confidence",
        "10_Adjusted_Valuation_QA",
        "11_Audit_Checks",
    ]
    assert workbook["10_Adjusted_Valuation_QA"]["B3"].value == 45000.0
    assert workbook["10_Adjusted_Valuation_QA"]["B4"].value == 80000.0
    assert workbook["06_DCF_Model"]["C2"].value.startswith("=B2*")

    presentation = Presentation(result.output_dir / "11_pitchbook.pptx")
    assert len(presentation.slides) == 10
    deck_text = "\n".join(
        shape.text for slide in presentation.slides for shape in slide.shapes if hasattr(shape, "text")
    )
    assert "DealForge AI" in deck_text
    assert "Adjusted Valuation QA" in deck_text
    assert "Human review required" in deck_text

    with ZipFile(result.bundle_path) as bundle:
        archived = set(bundle.namelist())
    assert archived == EXPECTED_ARTIFACTS


def test_text_outputs_do_not_leak_python_null_literals(tmp_path: Path) -> None:
    result = generate_portfolio_pack("sample_data", tmp_path / "pack")
    for name in [
        "01_company_profile.md",
        "02_investment_memo.md",
        "05_due_diligence_red_flags.md",
        "06_audit_trail.md",
        "16_investment_committee_memo.md",
        "17_due_diligence_request_list.md",
    ]:
        text = (result.output_dir / name).read_text(encoding="utf-8")
        assert "None" not in text
        assert "human review" in text.lower()
        assert "not investment advice" in text.lower()
        assert "23_adjusted_valuation_summary.json" in text or name in {
            "05_due_diligence_red_flags.md",
            "06_audit_trail.md",
        }


def test_portfolio_plus_json_artifacts_are_consistent(tmp_path: Path) -> None:
    result = generate_portfolio_pack("sample_data", tmp_path / "pack")
    source_index = json.loads((result.output_dir / "03_source_index.json").read_text(encoding="utf-8"))
    source_confidence = json.loads((result.output_dir / "18_source_confidence_ladder.json").read_text(encoding="utf-8"))
    audit = json.loads((result.output_dir / "19_valuation_input_audit.json").read_text(encoding="utf-8"))
    recommendations = json.loads((result.output_dir / "20_valuation_adjustment_recommendations.json").read_text(encoding="utf-8"))
    consistency = json.loads((result.output_dir / "15_consistency_checks.json").read_text(encoding="utf-8"))
    manifest = json.loads((result.output_dir / "00_deal_pack_bundle_manifest.json").read_text(encoding="utf-8"))

    assert source_index["document_count"] == 3
    assert source_confidence["status"] == "NEEDS_REVIEW"
    assert audit["status"] == "PASS_WITH_WARNINGS"
    assert recommendations["status"] == "PASS_WITH_RECOMMENDATIONS"
    assert recommendations["recommendation_count"] >= 10
    assert consistency["status"] == "PASS"
    assert len(manifest["files"]) == 23
    assert {row["file"] for row in manifest["files"]} == EXPECTED_ARTIFACTS - {"00_deal_pack_bundle_manifest.json"}
