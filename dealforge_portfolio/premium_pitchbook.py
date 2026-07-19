"""Premium PowerPoint upgrade layer for the public DealForge portfolio edition.

The core engine intentionally remains deterministic and synthetic. This module only
improves the generated pitchbook's presentation quality and then refreshes the
manifest/bundle so the downloadable pack stays internally consistent.
"""

from __future__ import annotations

import hashlib
import json
from datetime import datetime, timezone
from pathlib import Path
from typing import Any
from zipfile import ZIP_DEFLATED, ZipFile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


DARK = "07111F"
NAVY = "0B1F33"
TEAL = "0F766E"
GREEN = "22C55E"
AMBER = "F59E0B"
WHITE = "FFFFFF"
MUTED = "A7B0BE"
PANEL = "111827"
BORDER = "253144"


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _rgb(value: str) -> RGBColor:
    value = value.strip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _fmt(value: Any, suffix: str = "") -> str:
    if value is None:
        return "N/A"
    if isinstance(value, float):
        text = f"{value:,.2f}".rstrip("0").rstrip(".")
    else:
        text = str(value)
    return f"{text}{suffix}" if suffix else text


def _set_background(slide: Any, color: str = DARK) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _rect(slide: Any, left: float, top: float, width: float, height: float, fill: str, line: str | None = None, radius: bool = False) -> Any:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line:
        shape.line.color.rgb = _rgb(line)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def _text(
    slide: Any,
    value: str,
    left: float,
    top: float,
    width: float,
    height: float,
    size: int = 18,
    color: str = WHITE,
    bold: bool = False,
    align: Any | None = None,
) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    p.font.name = "Aptos"
    if align is not None:
        p.alignment = align
    return box


def _bullet_box(slide: Any, bullets: list[str], left: float, top: float, width: float, height: float, size: int = 15) -> Any:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(size)
        p.font.color.rgb = _rgb(WHITE)
        p.font.name = "Aptos"
        p.space_after = Pt(8)
    return box


def _header(slide: Any, title: str, subtitle: str = "") -> None:
    _rect(slide, 0, 0, 13.333, 0.72, NAVY)
    _text(slide, "DealForge AI", 0.45, 0.19, 2.4, 0.3, 13, "99F6E4", True)
    _text(slide, "Portfolio-plus public edition", 10.05, 0.19, 2.75, 0.3, 11, MUTED, False, PP_ALIGN.RIGHT)
    _text(slide, title, 0.65, 1.0, 8.8, 0.48, 28, WHITE, True)
    if subtitle:
        _text(slide, subtitle, 0.67, 1.52, 10.5, 0.34, 13, MUTED)


def _kpi(slide: Any, label: str, value: str, note: str, left: float, top: float, width: float = 2.55) -> None:
    _rect(slide, left, top, width, 1.15, PANEL, BORDER, radius=True)
    _text(slide, label.upper(), left + 0.18, top + 0.16, width - 0.34, 0.2, 8, MUTED, True)
    _text(slide, value, left + 0.18, top + 0.41, width - 0.34, 0.32, 18, WHITE, True)
    _text(slide, note, left + 0.18, top + 0.78, width - 0.34, 0.24, 8, MUTED)


def _status_badge(slide: Any, text: str, left: float, top: float, width: float = 2.15, fill: str = "3B2F08") -> None:
    _rect(slide, left, top, width, 0.38, fill, AMBER, radius=True)
    _text(slide, text, left + 0.08, top + 0.09, width - 0.16, 0.16, 8, "FDE68A", True, PP_ALIGN.CENTER)


def _table(slide: Any, rows: list[list[Any]], left: float, top: float, width: float, height: float, font_size: int = 10) -> None:
    if not rows:
        return
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(left), Inches(top), Inches(width), Inches(height)).table
    for col_idx in range(len(rows[0])):
        table.columns[col_idx].width = Inches(width / len(rows[0]))
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = _fmt(value)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb(NAVY if row_idx == 0 else DARK)
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = Pt(font_size)
            paragraph.font.name = "Aptos"
            paragraph.font.color.rgb = _rgb(WHITE if row_idx == 0 else "E5E7EB")
            paragraph.font.bold = row_idx == 0


def _new_slide(prs: Presentation, title: str, subtitle: str = "") -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _header(slide, title, subtitle)
    return slide


def write_premium_pitchbook(path: Path, summary: dict[str, Any]) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    ev_range = summary["adjusted_implied_enterprise_values"]["range"]
    audit = summary.get("valuation_input_audit", {})
    recommendations = summary.get("valuation_adjustment_recommendations", {})
    source_confidence = summary.get("source_confidence", {})
    peers = summary.get("peer_review_rows", [])
    precedents = summary.get("precedent_review_rows", [])
    metrics = summary.get("financial_extract", {}).get("metrics", [])

    # 1. Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, DARK)
    _rect(slide, 0, 0, 13.333, 7.5, DARK)
    _rect(slide, 0, 0, 13.333, 0.16, TEAL)
    _text(slide, "PORTFOLIO-PLUS PUBLIC EDITION", 0.7, 0.75, 5.2, 0.28, 12, "99F6E4", True)
    _text(slide, f"{summary['company']} Strategic Review", 0.7, 1.35, 9.3, 0.75, 34, WHITE, True)
    _text(slide, "DealForge AI — Investment Banking Analyst Copilot", 0.72, 2.15, 8.3, 0.38, 18, MUTED)
    _kpi(slide, "Adjusted EV range", f"{_fmt(ev_range['low'])}–{_fmt(ev_range['high'])}", summary["target_unit"], 0.75, 3.05, 3.0)
    _kpi(slide, "Deal-pack artifacts", str(summary.get("artifact_count", 24)), "Excel, PPT, memos, QA files", 3.95, 3.05, 3.0)
    _kpi(slide, "QA status", summary.get("adjusted_valuation_summary", {}).get("status", "PASS_WITH_RECOMMENDATIONS"), "review-ready scaffold", 7.15, 3.05, 3.8)
    _text(slide, "Synthetic portfolio demonstration only. Human review required. Not investment advice.", 0.72, 6.72, 9.0, 0.3, 11, MUTED)

    # 2. Executive snapshot
    slide = _new_slide(prs, "Executive Snapshot", "A repeatable finance-work-product compiler, not a generic chatbot.")
    _kpi(slide, "Revenue", _fmt(summary["target_revenue"]), summary["target_unit"], 0.65, 2.1)
    _kpi(slide, "EBITDA", _fmt(summary["target_ebitda"]), summary["target_unit"], 3.45, 2.1)
    _kpi(slide, "Peers kept", f"{summary['included_peer_count']} / {summary['included_peer_count'] + summary['excluded_peer_count']}", "target row removed", 6.25, 2.1)
    _kpi(slide, "Precedents kept", f"{summary['included_precedent_count']} / {summary['included_precedent_count'] + summary['excluded_precedent_count']}", "off-sector removed", 9.05, 2.1)
    _bullet_box(slide, [
        "The engine validates structured financial inputs, comparable companies, and precedent transactions.",
        "It preserves explicit inclusion/exclusion reasons instead of hiding them inside a median.",
        "The same adjusted valuation summary feeds Excel, PowerPoint, memos, diligence, JSON, CSVs, and ZIP manifest.",
    ], 0.8, 3.7, 11.2, 1.65, 15)

    # 3. Workflow map
    slide = _new_slide(prs, "Workflow Map", "Visible analyst controls from input sources through distributable work products.")
    steps = ["Source pack", "Citation map", "Financial extract", "Valuation audit", "Adjusted summary", "Excel / PPT / Memo", "ZIP bundle"]
    for idx, step in enumerate(steps):
        left = 0.55 + idx * 1.78
        _rect(slide, left, 2.2, 1.45, 0.8, PANEL, BORDER, radius=True)
        _text(slide, step, left + 0.1, 2.42, 1.25, 0.22, 10, WHITE, True, PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            _text(slide, "→", left + 1.5, 2.39, 0.24, 0.24, 18, "99F6E4", True, PP_ALIGN.CENTER)
    _bullet_box(slide, [
        "Artifacts are generated from one adjusted valuation spine to reduce cross-document mismatch.",
        "Human-review and not-investment-advice language is preserved across review workpapers.",
    ], 0.75, 4.0, 10.8, 1.1, 15)

    # 4. Source coverage
    slide = _new_slide(prs, "Source Coverage", "Synthetic sources are explicit and intentionally flagged as requiring replacement before real use.")
    docs = summary.get("source_index", {}).get("documents", [])
    rows = [["Document", "Type", "Description"]] + [[d.get("file_name"), d.get("source_type"), d.get("description")] for d in docs[:5]]
    _table(slide, rows, 0.65, 2.0, 11.9, 2.4, 10)
    _status_badge(slide, summary["source_verification_status"], 0.75, 5.2, 2.15)
    _text(slide, "Public-safe demo inputs keep the workflow inspectable without exposing confidential data rooms.", 3.1, 5.23, 8.7, 0.3, 12, MUTED)

    # 5. Financial snapshot
    slide = _new_slide(prs, "Financial Snapshot", "Selected source-backed metrics used by the valuation workflow.")
    metric_rows = [["Metric", "Period", "Value", "Unit", "Source chunk"]]
    for m in metrics[:6]:
        metric_rows.append([m.get("metric"), m.get("period"), m.get("value"), m.get("unit"), m.get("source_chunk_id")])
    _table(slide, metric_rows, 0.65, 2.0, 11.8, 2.7, 10)
    _bullet_box(slide, [
        "Metrics are deliberately structured in the public edition so recruiters can reproduce the pipeline quickly.",
        "Private product direction: document ingestion, AI extraction, and page-level citations remain Algosphere product scope.",
    ], 0.85, 5.05, 10.7, 1.1, 14)

    # 6. Comparable company QA
    slide = _new_slide(prs, "Comparable Company QA", "Target self-inclusion and NM EBITDA treatment are visible to the reviewer.")
    peer_rows = [["Company", "EV/Revenue", "EV/EBITDA", "Status", "Review notes"]]
    for row in peers[:5]:
        peer_rows.append([row.get("company"), row.get("ev_revenue_multiple"), row.get("ev_ebitda_multiple"), row.get("review_status"), row.get("review_notes")])
    _table(slide, peer_rows, 0.55, 1.95, 12.15, 3.05, 9)
    _kpi(slide, "Adjusted comps EV/Revenue", _fmt(summary["adjusted_medians"]["comps_ev_revenue"], "x"), "primary revenue multiple", 0.75, 5.45, 3.1)
    _kpi(slide, "Adjusted comps EV/EBITDA", _fmt(summary["adjusted_medians"]["comps_ev_ebitda"], "x"), "NM rows excluded", 4.1, 5.45, 3.1)

    # 7. Precedent QA
    slide = _new_slide(prs, "Precedent Transaction QA", "Sector mismatch and mixed transaction units are flagged before relying on precedent values.")
    precedent_rows = [["Target", "Acquirer", "EV/Revenue", "Unit", "Status", "Review notes"]]
    for row in precedents[:5]:
        precedent_rows.append([row.get("target"), row.get("acquirer"), row.get("ev_revenue_multiple"), row.get("transaction_value_unit"), row.get("review_status"), row.get("review_notes")])
    _table(slide, precedent_rows, 0.45, 1.95, 12.35, 3.05, 8)
    _kpi(slide, "Adjusted precedent EV/Revenue", _fmt(summary["adjusted_medians"]["precedent_ev_revenue"], "x"), "after exclusions", 0.75, 5.45, 3.3)
    _kpi(slide, "Unit normalization", "Needs review", summary["unit_normalization_status"], 4.35, 5.45, 3.3)

    # 8. Adjusted valuation QA
    slide = _new_slide(prs, "Adjusted Valuation QA", "One adjusted valuation spine powers all generated work products.")
    _kpi(slide, "Adjusted EV range low", _fmt(ev_range["low"]), summary["target_unit"], 0.75, 2.1, 3.0)
    _kpi(slide, "Adjusted EV range high", _fmt(ev_range["high"]), summary["target_unit"], 3.95, 2.1, 3.0)
    _kpi(slide, "Audit score", _fmt(audit.get("score")), f"{audit.get('warning_count', 0)} warning(s)", 7.15, 2.1, 3.0)
    _status_badge(slide, summary.get("adjusted_valuation_summary", {}).get("status", "PASS_WITH_RECOMMENDATIONS"), 0.85, 3.65, 3.1)
    _bullet_box(slide, [
        f"Adjusted comps EV/Revenue: {_fmt(summary['adjusted_medians']['comps_ev_revenue'], 'x')}",
        f"Adjusted comps EV/EBITDA: {_fmt(summary['adjusted_medians']['comps_ev_ebitda'], 'x')}",
        f"Adjusted precedent EV/Revenue: {_fmt(summary['adjusted_medians']['precedent_ev_revenue'], 'x')}",
        "Artifact spine: 23_adjusted_valuation_summary.json",
    ], 0.9, 4.25, 8.8, 1.6, 15)

    # 9. Review flags and diligence
    slide = _new_slide(prs, "Review Flags and Diligence", "Warnings are converted into deterministic analyst follow-ups.")
    _kpi(slide, "Audit warnings", str(audit.get("warning_count", 0)), "valuation input audit", 0.75, 2.05, 2.7)
    _kpi(slide, "Recommendations", str(recommendations.get("recommendation_count", 0)), "artifact 20", 3.75, 2.05, 2.7)
    _kpi(slide, "Source confidence", _fmt(source_confidence.get("top_score")), source_confidence.get("status", "NEEDS_REVIEW"), 6.75, 2.05, 2.7)
    rec_lines = [f"{r.get('recommendation_id')}: {r.get('item')} — {r.get('recommended_action')}" for r in recommendations.get("recommendations", [])[:5]]
    _bullet_box(slide, rec_lines or ["No review recommendations generated."], 0.85, 3.65, 11.4, 2.1, 12)

    # 10. Generated work products
    slide = _new_slide(prs, "Generated Work Products", "The portfolio edition proves end-to-end product thinking across model, memo, deck, and bundle outputs.")
    artifacts = [
        "10_valuation_model.xlsx",
        "11_pitchbook.pptx",
        "16_investment_committee_memo.md",
        "17_due_diligence_request_list.md",
        "20_valuation_adjustment_recommendations.json",
        "23_adjusted_valuation_summary.json",
        "00_deal_pack_bundle_manifest.json",
        "dealforge_portfolio_plus_pack.zip",
    ]
    rows = [["Artifact", "Purpose"]] + [[name, "Generated analyst work product"] for name in artifacts]
    _table(slide, rows, 0.75, 1.8, 11.6, 3.9, 10)
    _text(slide, "Human review required. Not investment advice, fairness opinion, or certified valuation.", 0.78, 6.45, 10.8, 0.35, 12, MUTED)

    prs.save(path)
    return path


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def refresh_manifest_and_bundle(result: Any) -> None:
    output_dir = Path(result.output_dir)
    manifest_path = output_dir / "00_deal_pack_bundle_manifest.json"
    bundle_path = Path(result.bundle_path)

    artifacts = [Path(path) for path in result.artifacts]
    artifact_names = [path.name for path in artifacts]
    manifest_inputs = [name for name in artifact_names if name != manifest_path.name]

    manifest = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "files": [
            {
                "file": name,
                "size_bytes": (output_dir / name).stat().st_size,
                "sha256": _sha256(output_dir / name),
            }
            for name in manifest_inputs
        ],
    }
    manifest_path.write_text(json.dumps(manifest, indent=2), encoding="utf-8")

    with ZipFile(bundle_path, "w", ZIP_DEFLATED) as bundle:
        for name in artifact_names:
            bundle.write(output_dir / name, arcname=name)


def upgrade_pitchbook_and_bundle(result: Any) -> None:
    write_premium_pitchbook(Path(result.output_dir) / "11_pitchbook.pptx", result.summary)
    refresh_manifest_and_bundle(result)
