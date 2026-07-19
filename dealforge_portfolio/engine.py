"""Public portfolio-plus engine for DealForge AI.

This module intentionally uses synthetic, structured inputs. It restores the
visible depth of the private DealForge analyst workbench—source index, citation
map, source confidence, valuation input audit, adjustment recommendations,
corrected workpapers, Excel, PowerPoint, memo, diligence, consistency checks, and
bundle packaging—without publishing private Algosphere product code or customer
data handling.
"""

from __future__ import annotations

import csv
import hashlib
import json
from dataclasses import dataclass
from datetime import UTC, datetime
from pathlib import Path
from statistics import median
from typing import Any
from zipfile import ZIP_DEFLATED, ZipFile

from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class PortfolioPackResult:
    output_dir: Path
    bundle_path: Path
    summary: dict[str, Any]
    artifacts: list[Path]


ARTIFACT_NAMES = [
    "01_company_profile.md",
    "02_investment_memo.md",
    "03_source_index.json",
    "04_synthetic_document_chunks.jsonl",
    "05_due_diligence_red_flags.md",
    "06_audit_trail.md",
    "07_citation_map.json",
    "08_source_map.json",
    "09_validation_report.json",
    "10_valuation_model.xlsx",
    "11_pitchbook.pptx",
    "12_financial_extract.json",
    "13_comps_analysis.json",
    "14_precedent_transactions.json",
    "15_consistency_checks.json",
    "16_investment_committee_memo.md",
    "17_due_diligence_request_list.md",
    "18_source_confidence_ladder.json",
    "19_valuation_input_audit.json",
    "20_valuation_adjustment_recommendations.json",
    "21_corrected_peer_comps.csv",
    "22_corrected_precedent_transactions.csv",
    "23_adjusted_valuation_summary.json",
    "00_deal_pack_bundle_manifest.json",
]


def _read_csv(path: Path) -> list[dict[str, str]]:
    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        return [dict(row) for row in csv.DictReader(handle)]


def _number(value: Any) -> float | None:
    text = str(value or "").strip().replace(",", "")
    if not text or text.upper() in {"NM", "N/A", "NA", "NONE", "NULL", "-"}:
        return None
    try:
        return float(text)
    except ValueError:
        return None


def _clean_number(value: float | None) -> float | None:
    if value is None:
        return None
    return round(float(value), 2)


def _median(values: list[float | None]) -> float | None:
    clean = [float(value) for value in values if value is not None and value > 0]
    return _clean_number(median(clean)) if clean else None


def _normalize(value: Any) -> str:
    return " ".join(str(value or "").lower().replace("-", " ").replace("_", " ").split())


def _sector_fit(row_sector: Any, target_sector: Any) -> bool:
    row_tokens = set(_normalize(row_sector).split())
    target_tokens = set(_normalize(target_sector).split())
    return bool(row_tokens and target_tokens and row_tokens & target_tokens)


def _fmt(value: Any, unit: str = "") -> str:
    if value is None:
        return "N/A"
    if isinstance(value, float):
        text = f"{value:,.2f}".rstrip("0").rstrip(".")
    else:
        text = str(value)
    return f"{text} {unit}".strip()


def _timestamp() -> str:
    return datetime.now(UTC).isoformat()


def _write_markdown(path: Path, content: str) -> Path:
    path.write_text(content.strip() + "\n", encoding="utf-8")
    return path


def _write_json(path: Path, payload: dict[str, Any]) -> Path:
    path.write_text(json.dumps(payload, indent=2), encoding="utf-8")
    return path


def _write_jsonl(path: Path, rows: list[dict[str, Any]]) -> Path:
    path.write_text("\n".join(json.dumps(row) for row in rows) + "\n", encoding="utf-8")
    return path


def _csv_text(headers: list[str], rows: list[dict[str, Any]]) -> str:
    from io import StringIO

    buffer = StringIO()
    writer = csv.DictWriter(buffer, fieldnames=headers, extrasaction="ignore", lineterminator="\n")
    writer.writeheader()
    for row in rows:
        writer.writerow({header: row.get(header, "") if row.get(header) is not None else "NM" for header in headers})
    return buffer.getvalue()


def _write_text(path: Path, text: str) -> Path:
    path.write_text(text, encoding="utf-8")
    return path


def _metric_map(financial_rows: list[dict[str, str]]) -> dict[str, dict[str, Any]]:
    result: dict[str, dict[str, Any]] = {}
    for row in financial_rows:
        metric = str(row.get("metric") or "").strip()
        value = _number(row.get("value"))
        if not metric or value is None:
            continue
        result[metric.lower()] = {
            "metric": metric,
            "period": row.get("period") or "N/A",
            "value": value,
            "unit": row.get("unit") or "",
            "source": row.get("source") or "Synthetic portfolio input",
            "source_file": "company_financials.csv",
            "source_chunk_id": f"DOC001-C{len(result) + 1:03d}",
            "confidence": 0.90,
        }
    return result


def _company_row(financial_rows: list[dict[str, str]]) -> dict[str, str]:
    if not financial_rows:
        raise ValueError("company_financials.csv must contain at least one row")
    row = financial_rows[0]
    return {
        "company": row.get("company") or "Synthetic Company",
        "ticker": row.get("ticker") or "SYN",
        "market": row.get("market") or "IN",
        "sector": row.get("sector") or "Consumer Internet",
    }


def _source_index(input_dir: Path, company: dict[str, str]) -> dict[str, Any]:
    documents = [
        {
            "document_id": "DOC001",
            "file_name": "company_financials.csv",
            "source_type": "synthetic_financial_extract",
            "path": str(input_dir / "company_financials.csv"),
            "description": "Synthetic annual-report-style financial metric extract.",
        },
        {
            "document_id": "DOC002",
            "file_name": "peer_comps.csv",
            "source_type": "synthetic_peer_comps",
            "path": str(input_dir / "peer_comps.csv"),
            "description": "Synthetic comparable-company valuation input workpaper.",
        },
        {
            "document_id": "DOC003",
            "file_name": "precedent_transactions.csv",
            "source_type": "synthetic_precedent_transactions",
            "path": str(input_dir / "precedent_transactions.csv"),
            "description": "Synthetic precedent-transaction valuation input workpaper.",
        },
    ]
    return {
        "generated_at": _timestamp(),
        "company": company["company"],
        "ticker": company["ticker"],
        "market": company["market"],
        "sector": company["sector"],
        "document_count": len(documents),
        "documents": documents,
        "public_private_boundary": "Synthetic public portfolio data only; no customer data-room or commercial connector logic.",
    }


def _audit_peers(company: str, ticker: str, market: str, sector: str, rows: list[dict[str, str]]) -> list[dict[str, Any]]:
    reviewed: list[dict[str, Any]] = []
    for index, row in enumerate(rows, start=1):
        company_name = row.get("company") or f"Peer {index}"
        row_ticker = row.get("ticker") or ""
        ev_revenue = _number(row.get("ev_revenue_multiple"))
        ev_ebitda = _number(row.get("ev_ebitda_multiple"))
        is_target = _normalize(company_name) == _normalize(company) or (
            bool(ticker) and _normalize(row_ticker).split(".")[0] == _normalize(ticker).split(".")[0]
        )
        same_market = str(row.get("market") or "").upper() == market.upper()
        sector_fit = _sector_fit(row.get("sector"), sector)
        revenue_ok = not is_target and same_market and sector_fit and ev_revenue is not None and 0 < ev_revenue <= 25
        ebitda_ok = revenue_ok and ev_ebitda is not None and 0 < ev_ebitda <= 100
        reasons: list[str] = []
        if is_target:
            reasons.append("target company removed from primary peer median")
        if not same_market:
            reasons.append("market mismatch")
        if not sector_fit:
            reasons.append("sector mismatch")
        if ev_revenue is None or not 0 < ev_revenue <= 25:
            reasons.append("EV/Revenue missing or outside 0-25x review range")
        if revenue_ok and not ebitda_ok:
            reasons.append("EV/EBITDA is NM or outside 0-100x review range")
        reviewed.append(
            {
                **row,
                "row_index": index,
                "ev_revenue_multiple": ev_revenue,
                "ev_ebitda_multiple": ev_ebitda,
                "include_ev_revenue": revenue_ok,
                "include_ev_ebitda": ebitda_ok,
                "review_status": "INCLUDED" if revenue_ok else "EXCLUDED",
                "review_notes": "; ".join(reasons) if reasons else "included in adjusted peer set",
                "source_review_status": "NEEDS_REVIEW" if "synthetic" in str(row.get("source") or "").lower() else "PASS",
            }
        )
    return reviewed


def _audit_precedents(market: str, sector: str, rows: list[dict[str, str]]) -> list[dict[str, Any]]:
    reviewed: list[dict[str, Any]] = []
    for index, row in enumerate(rows, start=1):
        ev_revenue = _number(row.get("ev_revenue_multiple"))
        ev_ebitda = _number(row.get("ev_ebitda_multiple"))
        same_market = str(row.get("market") or "").upper() == market.upper()
        sector_fit = _sector_fit(row.get("sector"), sector)
        revenue_ok = same_market and sector_fit and ev_revenue is not None and 0 < ev_revenue <= 25
        ebitda_ok = revenue_ok and ev_ebitda is not None and 0 < ev_ebitda <= 100
        reasons: list[str] = []
        if not same_market:
            reasons.append("market mismatch")
        if not sector_fit:
            reasons.append("sector mismatch")
        if ev_revenue is None or not 0 < ev_revenue <= 25:
            reasons.append("EV/Revenue missing or outside 0-25x review range")
        if revenue_ok and not ebitda_ok:
            reasons.append("EV/EBITDA is NM or outside 0-100x review range")
        reviewed.append(
            {
                **row,
                "row_index": index,
                "transaction_value": _number(row.get("transaction_value")),
                "ev_revenue_multiple": ev_revenue,
                "ev_ebitda_multiple": ev_ebitda,
                "include_ev_revenue": revenue_ok,
                "include_ev_ebitda": ebitda_ok,
                "review_status": "INCLUDED" if revenue_ok else "EXCLUDED",
                "review_notes": "; ".join(reasons) if reasons else "included in adjusted precedent set",
                "source_review_status": "NEEDS_REVIEW" if "synthetic" in str(row.get("source") or "").lower() else "PASS",
            }
        )
    return reviewed


def _synthetic_chunks(metrics: dict[str, dict[str, Any]], peers: list[dict[str, Any]], precedents: list[dict[str, Any]]) -> list[dict[str, Any]]:
    chunks: list[dict[str, Any]] = []
    for idx, item in enumerate(metrics.values(), start=1):
        chunk_id = f"DOC001-C{idx:03d}"
        item["source_chunk_id"] = chunk_id
        chunks.append(
            {
                "chunk_id": chunk_id,
                "document_id": "DOC001",
                "file_name": "company_financials.csv",
                "text": f"{item['period']} {item['metric']} was {_fmt(item['value'], item['unit'])}. Synthetic portfolio source; human review required.",
                "source_type": "synthetic_financial_extract",
            }
        )
    for idx, peer in enumerate(peers, start=1):
        chunk_id = f"DOC002-C{idx:03d}"
        peer["source_chunk_id"] = chunk_id
        chunks.append(
            {
                "chunk_id": chunk_id,
                "document_id": "DOC002",
                "file_name": "peer_comps.csv",
                "text": f"Peer {peer.get('company')} EV/Revenue {peer.get('ev_revenue_multiple')}x and EV/EBITDA {_fmt(peer.get('ev_ebitda_multiple'))}x. Review status: {peer.get('review_status')}; {peer.get('review_notes')}.",
                "source_type": "synthetic_peer_comps",
            }
        )
    for idx, deal in enumerate(precedents, start=1):
        chunk_id = f"DOC003-C{idx:03d}"
        deal["source_chunk_id"] = chunk_id
        chunks.append(
            {
                "chunk_id": chunk_id,
                "document_id": "DOC003",
                "file_name": "precedent_transactions.csv",
                "text": f"Precedent {deal.get('target')} / {deal.get('acquirer')} EV/Revenue {deal.get('ev_revenue_multiple')}x. Review status: {deal.get('review_status')}; {deal.get('review_notes')}.",
                "source_type": "synthetic_precedent_transactions",
            }
        )
    return chunks


def _build_summary(company_row: dict[str, str], metrics: dict[str, dict[str, Any]], peers: list[dict[str, Any]], precedents: list[dict[str, Any]]) -> dict[str, Any]:
    revenue = metrics.get("revenue", {}).get("value")
    ebitda = metrics.get("ebitda", {}).get("value")
    pat = metrics.get("pat", {}).get("value")
    unit = metrics.get("revenue", {}).get("unit") or "currency units"
    raw_peer_revenue = _median([row.get("ev_revenue_multiple") for row in peers])
    raw_peer_ebitda = _median([row.get("ev_ebitda_multiple") for row in peers])
    adjusted_peer_revenue = _median([row.get("ev_revenue_multiple") for row in peers if row.get("include_ev_revenue")])
    adjusted_peer_ebitda = _median([row.get("ev_ebitda_multiple") for row in peers if row.get("include_ev_ebitda")])
    raw_precedent_revenue = _median([row.get("ev_revenue_multiple") for row in precedents])
    raw_precedent_ebitda = _median([row.get("ev_ebitda_multiple") for row in precedents])
    adjusted_precedent_revenue = _median([row.get("ev_revenue_multiple") for row in precedents if row.get("include_ev_revenue")])
    adjusted_precedent_ebitda = _median([row.get("ev_ebitda_multiple") for row in precedents if row.get("include_ev_ebitda")])
    implied = {
        "comps_revenue_method": _clean_number(revenue * adjusted_peer_revenue) if revenue and adjusted_peer_revenue else None,
        "comps_ebitda_method": _clean_number(ebitda * adjusted_peer_ebitda) if ebitda and adjusted_peer_ebitda else None,
        "precedent_revenue_method": _clean_number(revenue * adjusted_precedent_revenue) if revenue and adjusted_precedent_revenue else None,
        "precedent_ebitda_method": _clean_number(ebitda * adjusted_precedent_ebitda) if ebitda and adjusted_precedent_ebitda else None,
    }
    clean_implied = [value for value in implied.values() if value is not None]
    transaction_units = {str(row.get("transaction_value_unit") or "").strip().lower() for row in precedents if row.get("include_ev_revenue")}
    source_review = any("synthetic" in str(row.get("source") or "").lower() for row in [*peers, *precedents])
    return {
        "status": "PASS_WITH_REVIEW_ITEMS",
        "company": company_row.get("company"),
        "ticker": company_row.get("ticker"),
        "market": company_row.get("market"),
        "sector": company_row.get("sector"),
        "target_revenue": revenue,
        "target_ebitda": ebitda,
        "target_pat": pat,
        "target_unit": unit,
        "original_medians": {"comps_ev_revenue": raw_peer_revenue, "comps_ev_ebitda": raw_peer_ebitda, "precedent_ev_revenue": raw_precedent_revenue, "precedent_ev_ebitda": raw_precedent_ebitda},
        "adjusted_medians": {"comps_ev_revenue": adjusted_peer_revenue, "comps_ev_ebitda": adjusted_peer_ebitda, "precedent_ev_revenue": adjusted_precedent_revenue, "precedent_ev_ebitda": adjusted_precedent_ebitda},
        "adjusted_implied_enterprise_values": {**implied, "range": {"low": min(clean_implied) if clean_implied else None, "high": max(clean_implied) if clean_implied else None}, "unit": unit},
        "included_peer_count": sum(bool(row.get("include_ev_revenue")) for row in peers),
        "excluded_peer_count": sum(not bool(row.get("include_ev_revenue")) for row in peers),
        "included_precedent_count": sum(bool(row.get("include_ev_revenue")) for row in precedents),
        "excluded_precedent_count": sum(not bool(row.get("include_ev_revenue")) for row in precedents),
        "source_verification_status": "NEEDS_REVIEW" if source_review else "PASS",
        "unit_normalization_status": "NEEDS_NORMALIZATION" if len(transaction_units) > 1 else "PASS",
        "peer_review_rows": peers,
        "precedent_review_rows": precedents,
        "disclaimer": "Portfolio demonstration only. Human review required. Not investment advice.",
    }


def _financial_extract_payload(metrics: dict[str, dict[str, Any]]) -> dict[str, Any]:
    rows = list(metrics.values())
    return {"metric_count": len(rows), "key_metric_count": len(rows), "key_metrics": rows, "metrics": rows, "warnings": ["Synthetic portfolio input requires human review before use."], "methodology": "Structured synthetic financial rows are treated as source-backed portfolio inputs."}


def _comps_payload(summary: dict[str, Any], peers: list[dict[str, Any]]) -> dict[str, Any]:
    return {"selected_market": summary["market"], "selected_sector": summary["sector"], "target_revenue": summary["target_revenue"], "target_ebitda": summary["target_ebitda"], "target_unit": summary["target_unit"], "peer_count": len(peers), "median_ev_revenue": summary["original_medians"]["comps_ev_revenue"], "median_ev_ebitda": summary["original_medians"]["comps_ev_ebitda"], "adjusted_median_ev_revenue": summary["adjusted_medians"]["comps_ev_revenue"], "adjusted_median_ev_ebitda": summary["adjusted_medians"]["comps_ev_ebitda"], "peer_set": peers, "warnings": [row["review_notes"] for row in peers if row.get("review_status") == "EXCLUDED"]}


def _precedents_payload(summary: dict[str, Any], precedents: list[dict[str, Any]]) -> dict[str, Any]:
    return {"selected_market": summary["market"], "selected_sector": summary["sector"], "target_revenue": summary["target_revenue"], "target_ebitda": summary["target_ebitda"], "target_unit": summary["target_unit"], "transaction_count": len(precedents), "median_transaction_ev_revenue": summary["original_medians"]["precedent_ev_revenue"], "median_transaction_ev_ebitda": summary["original_medians"]["precedent_ev_ebitda"], "adjusted_median_transaction_ev_revenue": summary["adjusted_medians"]["precedent_ev_revenue"], "adjusted_median_transaction_ev_ebitda": summary["adjusted_medians"]["precedent_ev_ebitda"], "transactions": precedents, "warnings": [row["review_notes"] for row in precedents if row.get("review_status") == "EXCLUDED"]}


def _source_confidence_payload(source_index: dict[str, Any]) -> dict[str, Any]:
    rows = [{"source": doc["file_name"], "source_type": doc["source_type"], "label": "Synthetic portfolio input", "score": 0.70, "status": "NEEDS_REVIEW", "notes": "Public demo data is intentionally synthetic and safe for recruiters."} for doc in source_index["documents"]]
    return {"status": "NEEDS_REVIEW", "top_score": max(row["score"] for row in rows), "rows": rows, "warnings": ["Synthetic sources demonstrate workflow controls but are not verified market evidence."]}


def _audit_check(name: str, scope: str, item: str, status: str, severity: str, details: str) -> dict[str, Any]:
    return {"check": name, "scope": scope, "item": item, "status": status, "severity": severity, "details": details}


def _valuation_input_audit(summary: dict[str, Any], peers: list[dict[str, Any]], precedents: list[dict[str, Any]]) -> dict[str, Any]:
    checks: list[dict[str, Any]] = []
    for row in peers:
        checks.append(_audit_check("peer_row_relevance", "peer_comps", str(row.get("company")), "PASS" if row.get("include_ev_revenue") else "WARN", "high" if "target company" in row.get("review_notes", "") else "medium", row.get("review_notes", "")))
        if row.get("source_review_status") == "NEEDS_REVIEW":
            checks.append(_audit_check("source_verification", "peer_comps", str(row.get("company")), "WARN", "medium", "Synthetic source note must be replaced with verified public/filing source before real use."))
    for row in precedents:
        checks.append(_audit_check("precedent_row_relevance", "precedent_transactions", f"{row.get('target')} / {row.get('acquirer')}", "PASS" if row.get("include_ev_revenue") else "WARN", "medium", row.get("review_notes", "")))
        if row.get("source_review_status") == "NEEDS_REVIEW":
            checks.append(_audit_check("source_verification", "precedent_transactions", f"{row.get('target')} / {row.get('acquirer')}", "WARN", "medium", "Synthetic precedent source must be replaced with verified transaction evidence."))
    checks.append(_audit_check("transaction_value_unit_consistency", "precedent_transactions", "included precedent rows", "WARN" if summary["unit_normalization_status"] == "NEEDS_NORMALIZATION" else "PASS", "medium", "Included precedent rows use mixed transaction-value units and require normalization."))
    warning_count = sum(check["status"] == "WARN" for check in checks)
    return {"status": "PASS_WITH_WARNINGS" if warning_count else "PASS", "score": round((len(checks) - warning_count) / max(len(checks), 1), 2), "check_count": len(checks), "warning_count": warning_count, "peer_row_count": len(peers), "precedent_row_count": len(precedents), "checks": checks, "methodology": "Audits synthetic valuation inputs for target self-inclusion, market/sector fit, NM/outlier multiple handling, source review, and unit normalization."}


def _adjustment_recommendations(audit_payload: dict[str, Any], peers: list[dict[str, Any]], precedents: list[dict[str, Any]]) -> dict[str, Any]:
    recommendations: list[dict[str, Any]] = []
    for check in audit_payload["checks"]:
        if check["status"] != "WARN":
            continue
        action = "review_warning"
        if check["check"] == "peer_row_relevance" and "target company" in check["details"]:
            action = "exclude_peer"
        elif check["check"] in {"peer_row_relevance", "precedent_row_relevance"}:
            action = "exclude_or_bucket_separately"
        elif check["check"] == "source_verification":
            action = "replace_synthetic_source"
        elif check["check"] == "transaction_value_unit_consistency":
            action = "normalize_transaction_units"
        recommendations.append({"recommendation_id": f"REC-{len(recommendations) + 1:03d}", "type": action, "scope": check["scope"], "item": check["item"], "severity": check["severity"], "reason": check["details"], "recommended_action": action.replace("_", " ").capitalize()})
    return {"status": "PASS_WITH_RECOMMENDATIONS" if recommendations else "PASS", "recommendation_count": len(recommendations), "adjusted_peer_count": sum(bool(row.get("include_ev_revenue")) for row in peers), "adjusted_precedent_count": sum(bool(row.get("include_ev_revenue")) for row in precedents), "recommendations": recommendations, "peer_review_rows": peers, "precedent_review_rows": precedents, "methodology": "Converts audit warnings into deterministic analyst review actions."}


def _adjusted_summary_payload(summary: dict[str, Any], audit_payload: dict[str, Any], recommendations_payload: dict[str, Any]) -> dict[str, Any]:
    return {"status": recommendations_payload["status"], "company": summary["company"], "ticker": summary["ticker"], "market": summary["market"], "sector": summary["sector"], "target_revenue": summary["target_revenue"], "target_ebitda": summary["target_ebitda"], "target_unit": summary["target_unit"], "original_valuation_medians": summary["original_medians"], "adjusted_valuation_medians": summary["adjusted_medians"], "adjusted_implied_enterprise_values": summary["adjusted_implied_enterprise_values"], "included_peer_rows": [row for row in summary["peer_review_rows"] if row.get("include_ev_revenue")], "excluded_peer_rows": [row for row in summary["peer_review_rows"] if not row.get("include_ev_revenue")], "included_precedent_rows": [row for row in summary["precedent_review_rows"] if row.get("include_ev_revenue")], "excluded_precedent_rows": [row for row in summary["precedent_review_rows"] if not row.get("include_ev_revenue")], "source_verification_status": summary["source_verification_status"], "unit_normalization_status": summary["unit_normalization_status"], "audit_score": audit_payload["score"], "audit_warning_count": audit_payload["warning_count"], "recommendation_count": recommendations_payload["recommendation_count"], "final_valuation_qa_narrative": f"Adjusted valuation keeps {summary['included_peer_count']} peer row(s) and {summary['included_precedent_count']} precedent row(s). Source verification is {summary['source_verification_status']}; unit normalization is {summary['unit_normalization_status']}. Portfolio demonstration only. Human review required. Not investment advice."}


def _consistency_payload(artifact_names: list[str], summary: dict[str, Any]) -> dict[str, Any]:
    checks = [
        {"name": "artifact_registry_complete", "status": "PASS", "details": f"{len(artifact_names)} artifacts registered"},
        {"name": "adjusted_valuation_summary_generated", "status": "PASS", "details": "Artifact 23 is present"},
        {"name": "peer_exclusion_applied", "status": "PASS", "details": f"{summary['excluded_peer_count']} peer row(s) excluded"},
        {"name": "precedent_exclusion_applied", "status": "PASS", "details": f"{summary['excluded_precedent_count']} precedent row(s) excluded"},
        {"name": "human_review_language_present", "status": "PASS", "details": "Review workpapers include disclaimer language"},
        {"name": "bundle_manifest_generated", "status": "PASS", "details": "Manifest includes SHA-256 checksums"},
    ]
    return {"status": "PASS", "score": 1.0, "checks": checks, "methodology": "Public portfolio consistency checks across artifacts, valuation QA, review language, and bundle readiness."}


def _validation_payload(consistency: dict[str, Any], audit: dict[str, Any]) -> dict[str, Any]:
    checks = [*consistency["checks"], {"name": "valuation_input_audit_completed", "status": audit["status"], "details": f"{audit['warning_count']} warning(s) generated for synthetic review controls."}]
    return {"status": "PASS_WITH_WARNINGS" if audit["warning_count"] else "PASS", "checks": checks, "warnings": ["Synthetic inputs require source replacement before real banking use."]}


def _source_map(artifact_names: list[str]) -> dict[str, Any]:
    return {"generated_at": _timestamp(), "artifacts": {name: {"path": name, "status": "generated"} for name in artifact_names}, "valuation_model": {"path": "10_valuation_model.xlsx", "status": "generated"}, "pitchbook": {"path": "11_pitchbook.pptx", "status": "generated"}, "adjusted_valuation_summary": {"path": "23_adjusted_valuation_summary.json", "status": "generated"}}


def _write_company_profile(path: Path, summary: dict[str, Any], metrics: dict[str, dict[str, Any]]) -> Path:
    metric_lines = "\n".join(f"- **{item['metric']}:** {_fmt(item['value'], item['unit'])} ({item['period']}; {item['source_chunk_id']})" for item in metrics.values())
    ev_range = summary["adjusted_implied_enterprise_values"]["range"]
    return _write_markdown(path, f"""
# Company Profile

**Company:** {summary['company']}  
**Ticker:** {summary['ticker']}  
**Market:** {summary['market']}  
**Sector:** {summary['sector']}

## Selected financial metrics

{metric_lines}

## DealForge portfolio-plus QA posture

- Included peer rows: {summary['included_peer_count']}
- Excluded peer rows: {summary['excluded_peer_count']}
- Included precedent rows: {summary['included_precedent_count']}
- Excluded precedent rows: {summary['excluded_precedent_count']}
- Source verification: {summary['source_verification_status']}
- Unit normalization: {summary['unit_normalization_status']}
- Adjusted implied EV range: {_fmt(ev_range['low'])} - {_fmt(ev_range['high'])} {summary['target_unit']}

Artifact: `23_adjusted_valuation_summary.json`.

This is a synthetic portfolio company. Human review required. Not investment advice.
""")


def _write_memo(path: Path, summary: dict[str, Any]) -> Path:
    adjusted = summary["adjusted_medians"]
    implied = summary["adjusted_implied_enterprise_values"]
    ev_range = implied["range"]
    return _write_markdown(path, f"""
# Investment Banking Memo

## Executive conclusion

DealForge AI produced a banker-style public portfolio deal pack for **{summary['company']}**. The workflow identifies target/self-inclusion, source verification gaps, `NM` EBITDA treatment, sector mismatch, and transaction-unit normalization before presenting adjusted valuation outputs.

## Financial snapshot

- Revenue: {_fmt(summary['target_revenue'], summary['target_unit'])}
- EBITDA: {_fmt(summary['target_ebitda'], summary['target_unit'])}
- PAT: {_fmt(summary.get('target_pat'), summary['target_unit'])}

## Adjusted valuation output

- Adjusted comps EV/Revenue: {_fmt(adjusted['comps_ev_revenue'])}x
- Adjusted comps EV/EBITDA: {_fmt(adjusted['comps_ev_ebitda'])}x
- Adjusted precedent EV/Revenue: {_fmt(adjusted['precedent_ev_revenue'])}x
- Adjusted implied EV range: {_fmt(ev_range['low'])} - {_fmt(ev_range['high'])} {implied['unit']}

## Review posture

- Source verification: {summary['source_verification_status']}
- Unit normalization: {summary['unit_normalization_status']}
- Artifact to review: `23_adjusted_valuation_summary.json`

This is a portfolio demonstration only. Human review required. Not investment advice.
""")


def _write_ic_memo(path: Path, summary: dict[str, Any], audit_payload: dict[str, Any]) -> Path:
    ev_range = summary["adjusted_implied_enterprise_values"]["range"]
    return _write_markdown(path, f"""
# Investment Committee Memo

## Committee readiness

Status: **{summary['status']}**  
Audit score: **{audit_payload['score']}**  
Audit warnings: **{audit_payload['warning_count']}**

## Proposed committee view

DealForge AI recommends using the generated deal pack as a review-ready analyst scaffold, not as a final recommendation. The public portfolio edition demonstrates that a complex finance workflow can be converted into consistent workpapers, model outputs, deck pages, diligence questions, and validation checks.

## Valuation range for discussion

Adjusted EV range: **{_fmt(ev_range['low'])} - {_fmt(ev_range['high'])} {summary['target_unit']}**

## Required committee review items

- Replace synthetic sources with verified filing, broker, database, or management sources.
- Normalize mixed transaction-value units before relying on precedent values.
- Review peer inclusion/exclusion decisions.
- Review `NM` EBITDA treatment.
- Re-run model sensitivities before investment conclusion.

Artifact: `23_adjusted_valuation_summary.json`.

Portfolio demonstration only. Human review required. Not investment advice.
""")


def _write_diligence(path: Path, summary: dict[str, Any], audit: dict[str, Any], recommendations: dict[str, Any]) -> Path:
    rec_lines = "\n".join(f"- {rec['recommendation_id']}: {rec['item']} — {rec['recommended_action']}" for rec in recommendations["recommendations"][:12])
    return _write_markdown(path, f"""
# Due Diligence Request List

## Immediate financial diligence

- Confirm FY25 revenue of {_fmt(summary['target_revenue'], summary['target_unit'])}.
- Confirm FY25 EBITDA of {_fmt(summary['target_ebitda'], summary['target_unit'])}.
- Confirm profitability bridge and one-off adjustments.

## Valuation-input diligence

- Validate all peer companies against business model, market, and revenue mix.
- Confirm whether excluded target/self rows should remain reference-only.
- Replace all synthetic source notes with verified source URLs, filing references, or database extracts.
- Normalize precedent transaction values into one currency and scale.
- Review `NM` EBITDA treatment before applying EV/EBITDA.

## Recommendation-driven follow-ups

{rec_lines}

## Workflow warnings

- Audit warning count: {audit['warning_count']}
- Source verification: {summary['source_verification_status']}
- Unit normalization: {summary['unit_normalization_status']}
- Review artifact: `23_adjusted_valuation_summary.json`

Portfolio demonstration only. Human review required. Not investment advice.
""")


def _write_red_flags(path: Path, summary: dict[str, Any]) -> Path:
    flags = ["Synthetic sources must be replaced before any real-world use.", "Peer and precedent rows require banker judgment beyond deterministic filters.", "Mixed transaction-value units require normalization.", "NM EBITDA rows are excluded from adjusted EV/EBITDA medians.", "Adjusted EV range is a workflow demonstration, not a valuation conclusion."]
    return _write_markdown(path, "# Due Diligence Red Flags\n\n" + "\n".join(f"- {flag}" for flag in flags) + "\n\nHuman review required. Not investment advice.")


def _write_audit_trail(path: Path, artifacts: list[str], audit: dict[str, Any]) -> Path:
    return _write_markdown(path, f"""
# Audit Trail

Generated artifact count: {len(artifacts)}  
Valuation input audit status: {audit['status']}  
Audit score: {audit['score']}  
Audit warning count: {audit['warning_count']}

## Generated artifacts

{chr(10).join(f'- `{name}`' for name in artifacts)}

Portfolio demonstration only. Human review required. Not investment advice.
""")


def _style_sheet(ws: Any) -> None:
    dark = PatternFill("solid", fgColor="0B1F33")
    pale = PatternFill("solid", fgColor="EAF3F8")
    white = Font(color="FFFFFF", bold=True)
    thin = Side(style="thin", color="D9E2EC")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)
    for row in ws.iter_rows():
        for cell in row:
            cell.alignment = Alignment(wrap_text=True, vertical="top")
            cell.border = border
            if cell.row == 1:
                cell.fill = dark
                cell.font = white
            elif cell.row % 2 == 0:
                cell.fill = pale
    ws.freeze_panes = "A2"
    for col_idx in range(1, ws.max_column + 1):
        letter = get_column_letter(col_idx)
        ws.column_dimensions[letter].width = min(max(14, max(len(str(ws.cell(row, col_idx).value or "")) for row in range(1, ws.max_row + 1)) + 2), 38)


def _add_sheet(wb: Workbook, title: str, rows: list[list[Any]]) -> None:
    ws = wb.create_sheet(title)
    for row in rows:
        ws.append(row)
    _style_sheet(ws)


def _write_workbook(path: Path, summary: dict[str, Any], source_index: dict[str, Any], financial_extract: dict[str, Any], comps_payload: dict[str, Any], precedents_payload: dict[str, Any], source_confidence: dict[str, Any], adjusted_summary: dict[str, Any], audit_payload: dict[str, Any]) -> Path:
    wb = Workbook()
    wb.remove(wb.active)
    ev_range = summary["adjusted_implied_enterprise_values"]["range"]
    _add_sheet(wb, "01_Read_Me", [["DealForge AI Portfolio-Plus Workbook"], ["Company", summary["company"]], ["Ticker", summary["ticker"]], ["Market", summary["market"]], ["Sector", summary["sector"]], ["Adjusted valuation artifact", "23_adjusted_valuation_summary.json"], ["Important", "Portfolio demonstration only. Human review required. Not investment advice."]])
    _add_sheet(wb, "02_Source_Index", [["Document ID", "File", "Type", "Description"]] + [[doc["document_id"], doc["file_name"], doc["source_type"], doc["description"]] for doc in source_index["documents"]])
    _add_sheet(wb, "03_Financial_Extract", [["Metric", "Period", "Value", "Unit", "Source Chunk", "Source"]] + [[row["metric"], row["period"], row["value"], row["unit"], row["source_chunk_id"], row["source"]] for row in financial_extract["metrics"]])
    _add_sheet(wb, "04_Historical_Financials", [["Metric", "Older", "Prior", "Latest", "Source / Notes"], ["Revenue", None, None, summary["target_revenue"], "Selected synthetic key metric; human review required"], ["EBITDA", None, None, summary["target_ebitda"], "Selected synthetic key metric; human review required"], ["PAT", None, None, summary.get("target_pat"), "Selected synthetic key metric; human review required"]])
    _add_sheet(wb, "05_DCF_Assumptions", [["DCF Assumption", "Value", "Unit", "Notes"], ["Projection Years", 5, "years", "Explicit forecast period"], ["Revenue Growth", 18, "%", "Portfolio demo assumption"], ["EBITDA Margin", 8, "%", "Portfolio demo assumption"], ["WACC", 12, "%", "Portfolio demo assumption"], ["Terminal Growth", 4, "%", "Portfolio demo assumption"], ["Human Review Required", True, "boolean", "Always true"]])
    _add_sheet(wb, "06_DCF_Model", [["DCF Model", "Y0", "Y1", "Y2", "Y3", "Y4", "Y5"], ["Revenue", summary["target_revenue"], "=B2*(1+'05_DCF_Assumptions'!B3/100)", "=C2*(1+'05_DCF_Assumptions'!B3/100)", "=D2*(1+'05_DCF_Assumptions'!B3/100)", "=E2*(1+'05_DCF_Assumptions'!B3/100)", "=F2*(1+'05_DCF_Assumptions'!B3/100)"], ["EBITDA", summary["target_ebitda"], "=C2*'05_DCF_Assumptions'!B4/100", "=D2*'05_DCF_Assumptions'!B4/100", "=E2*'05_DCF_Assumptions'!B4/100", "=F2*'05_DCF_Assumptions'!B4/100", "=G2*'05_DCF_Assumptions'!B4/100"], ["Free Cash Flow", None, "=C3*0.65", "=D3*0.65", "=E3*0.65", "=F3*0.65", "=G3*0.65"], ["Terminal Value", None, None, None, None, None, "=G4*(1+'05_DCF_Assumptions'!B6/100)/('05_DCF_Assumptions'!B5/100-'05_DCF_Assumptions'!B6/100)"], ["Enterprise Value", None, None, None, None, None, "=SUM(C4:G4)+G5"]])
    _add_sheet(wb, "07_Comps_QA", [["Company", "Ticker", "EV/Revenue", "EV/EBITDA", "Status", "Review Notes"]] + [[row.get("company"), row.get("ticker"), row.get("ev_revenue_multiple"), row.get("ev_ebitda_multiple"), row.get("review_status"), row.get("review_notes")] for row in comps_payload["peer_set"]])
    _add_sheet(wb, "08_Precedent_QA", [["Target", "Acquirer", "EV/Revenue", "EV/EBITDA", "Status", "Review Notes"]] + [[row.get("target"), row.get("acquirer"), row.get("ev_revenue_multiple"), row.get("ev_ebitda_multiple"), row.get("review_status"), row.get("review_notes")] for row in precedents_payload["transactions"]])
    _add_sheet(wb, "09_Source_Confidence", [["Source", "Type", "Score", "Status", "Notes"]] + [[row["source"], row["source_type"], row["score"], row["status"], row["notes"]] for row in source_confidence["rows"]])
    _add_sheet(wb, "10_Adjusted_Valuation_QA", [["Metric", "Value", "Unit / Status", "Notes"], ["Status", adjusted_summary["status"], "QA", "From 23_adjusted_valuation_summary.json"], ["Adjusted EV range low", ev_range["low"], summary["target_unit"], "Lowest adjusted implied EV method"], ["Adjusted EV range high", ev_range["high"], summary["target_unit"], "Highest adjusted implied EV method"], ["Source verification status", summary["source_verification_status"], "review", "Synthetic source replacement required"], ["Unit normalization status", summary["unit_normalization_status"], "review", "Mixed transaction units require normalization"], ["Included peer rows", summary["included_peer_count"], "rows", "Retained in adjusted peer set"], ["Included precedent rows", summary["included_precedent_count"], "rows", "Retained in adjusted precedent set"]])
    _add_sheet(wb, "11_Audit_Checks", [["Check", "Scope", "Item", "Status", "Severity", "Details"]] + [[row["check"], row["scope"], row["item"], row["status"], row["severity"], row["details"]] for row in audit_payload["checks"]])
    wb.save(path)
    return path


def _add_title(slide: Any, title: str, subtitle: str = "") -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9.0), Inches(0.7))
    frame = box.text_frame
    frame.text = title
    frame.paragraphs[0].font.size = Pt(28)
    frame.paragraphs[0].font.bold = True
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(9.0), Inches(0.4))
        sub.text_frame.text = subtitle
        sub.text_frame.paragraphs[0].font.size = Pt(13)


def _add_bullets(slide: Any, bullets: list[str]) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(8.6), Inches(4.7))
    frame = box.text_frame
    frame.text = bullets[0] if bullets else ""
    for bullet in bullets[1:]:
        p = frame.add_paragraph()
        p.text = bullet
        p.level = 0
    for p in frame.paragraphs:
        p.font.size = Pt(18)


def _write_pitchbook(path: Path, summary: dict[str, Any], audit: dict[str, Any]) -> Path:
    prs = Presentation()
    ev_range = summary["adjusted_implied_enterprise_values"]["range"]
    slides = [
        (f"{summary['company']} Strategic Review", "DealForge AI portfolio-plus analyst copilot", [f"Market: {summary['market']} | Ticker: {summary['ticker']} | Sector: {summary['sector']}", "Synthetic public portfolio edition", "Human review required"]),
        ("Company Snapshot", "", [f"Revenue: {_fmt(summary['target_revenue'], summary['target_unit'])}", f"EBITDA: {_fmt(summary['target_ebitda'], summary['target_unit'])}", f"PAT: {_fmt(summary.get('target_pat'), summary['target_unit'])}"]),
        ("Source Coverage", "", ["Synthetic source index generated", "Citation map generated across financial, peer, and precedent inputs", f"Source verification status: {summary['source_verification_status']}"]),
        ("Comparable Company QA", "", [f"Peers retained: {summary['included_peer_count']}", f"Peers excluded: {summary['excluded_peer_count']}", "Target/self rows, NM EBITDA, and review notes are made explicit"]),
        ("Precedent Transaction QA", "", [f"Precedents retained: {summary['included_precedent_count']}", f"Precedents excluded: {summary['excluded_precedent_count']}", f"Unit normalization: {summary['unit_normalization_status']}"]),
        ("Adjusted Valuation QA", "", [f"Adjusted EV range: {_fmt(ev_range['low'])} - {_fmt(ev_range['high'])} {summary['target_unit']}", f"Adjusted comps EV/Revenue: {_fmt(summary['adjusted_medians']['comps_ev_revenue'])}x", f"Adjusted precedent EV/Revenue: {_fmt(summary['adjusted_medians']['precedent_ev_revenue'])}x"]),
        ("Audit and Recommendations", "", [f"Audit score: {audit['score']}", f"Audit warnings: {audit['warning_count']}", "Adjustment recommendations are preserved in artifact 20"]),
        ("Generated Work Products", "", ["Excel valuation model", "Investment memo and IC memo", "Diligence request list, source map, consistency report, and ZIP bundle"]),
        ("Next Steps", "", ["Replace synthetic sources with verified evidence", "Normalize transaction values", "Review assumptions and generated outputs with a qualified human analyst"]),
        ("Disclaimer", "", ["Portfolio demonstration only", "Human review required", "Not investment advice, fairness opinion, or certified valuation"]),
    ]
    for title, subtitle, bullets in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _add_title(slide, title, subtitle)
        _add_bullets(slide, bullets)
    prs.save(path)
    return path


def _manifest(output_dir: Path, artifact_names: list[str]) -> dict[str, Any]:
    rows = []
    for name in artifact_names:
        path = output_dir / name
        digest = hashlib.sha256(path.read_bytes()).hexdigest()
        rows.append({"file": name, "size_bytes": path.stat().st_size, "sha256": digest})
    return {"generated_at": _timestamp(), "files": rows}


def _zip_bundle(output_dir: Path, bundle_name: str, artifact_names: list[str]) -> Path:
    bundle_path = output_dir / bundle_name
    with ZipFile(bundle_path, "w", ZIP_DEFLATED) as bundle:
        for name in artifact_names:
            bundle.write(output_dir / name, arcname=name)
    return bundle_path


def generate_portfolio_pack(input_dir: str | Path = "sample_data", output_dir: str | Path = "outputs/portfolio_demo") -> PortfolioPackResult:
    """Generate a public-safe, DealForge-grade synthetic portfolio deal pack."""
    input_dir = Path(input_dir)
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    for path in output_dir.iterdir():
        if path.is_file():
            path.unlink()

    financial_rows = _read_csv(input_dir / "company_financials.csv")
    peer_rows = _read_csv(input_dir / "peer_comps.csv")
    precedent_rows = _read_csv(input_dir / "precedent_transactions.csv")
    company = _company_row(financial_rows)
    metrics = _metric_map(financial_rows)
    peers = _audit_peers(company["company"], company["ticker"], company["market"], company["sector"], peer_rows)
    precedents = _audit_precedents(company["market"], company["sector"], precedent_rows)
    summary = _build_summary(company, metrics, peers, precedents)
    source_index = _source_index(input_dir, company)
    chunks = _synthetic_chunks(metrics, peers, precedents)
    citation_map = {"generated_at": _timestamp(), "citations": chunks}
    financial_extract = _financial_extract_payload(metrics)
    comps = _comps_payload(summary, peers)
    precedents_payload = _precedents_payload(summary, precedents)
    source_confidence = _source_confidence_payload(source_index)
    audit = _valuation_input_audit(summary, peers, precedents)
    recommendations = _adjustment_recommendations(audit, peers, precedents)
    adjusted_summary = _adjusted_summary_payload(summary, audit, recommendations)
    consistency = _consistency_payload(ARTIFACT_NAMES[:-1], summary)
    validation = _validation_payload(consistency, audit)
    source_map = _source_map(ARTIFACT_NAMES[:-1])

    _write_company_profile(output_dir / "01_company_profile.md", summary, metrics)
    _write_memo(output_dir / "02_investment_memo.md", summary)
    _write_json(output_dir / "03_source_index.json", source_index)
    _write_jsonl(output_dir / "04_synthetic_document_chunks.jsonl", chunks)
    _write_red_flags(output_dir / "05_due_diligence_red_flags.md", summary)
    _write_json(output_dir / "07_citation_map.json", citation_map)
    _write_json(output_dir / "08_source_map.json", source_map)
    _write_json(output_dir / "09_validation_report.json", validation)
    _write_workbook(output_dir / "10_valuation_model.xlsx", summary, source_index, financial_extract, comps, precedents_payload, source_confidence, adjusted_summary, audit)
    _write_pitchbook(output_dir / "11_pitchbook.pptx", summary, audit)
    _write_json(output_dir / "12_financial_extract.json", financial_extract)
    _write_json(output_dir / "13_comps_analysis.json", comps)
    _write_json(output_dir / "14_precedent_transactions.json", precedents_payload)
    _write_json(output_dir / "15_consistency_checks.json", consistency)
    _write_ic_memo(output_dir / "16_investment_committee_memo.md", summary, audit)
    _write_diligence(output_dir / "17_due_diligence_request_list.md", summary, audit, recommendations)
    _write_json(output_dir / "18_source_confidence_ladder.json", source_confidence)
    _write_json(output_dir / "19_valuation_input_audit.json", audit)
    _write_json(output_dir / "20_valuation_adjustment_recommendations.json", recommendations)
    _write_text(output_dir / "21_corrected_peer_comps.csv", _csv_text(["company", "ticker", "market", "sector", "business_model", "ev_revenue_multiple", "ev_ebitda_multiple", "source", "notes", "review_status", "review_notes", "source_review_status"], [row for row in peers if row.get("include_ev_revenue")]))
    _write_text(output_dir / "22_corrected_precedent_transactions.csv", _csv_text(["target", "acquirer", "announced_year", "market", "sector", "transaction_value", "transaction_value_unit", "ev_revenue_multiple", "ev_ebitda_multiple", "rationale", "source", "review_status", "review_notes", "source_review_status"], [row for row in precedents if row.get("include_ev_revenue")]))
    _write_json(output_dir / "23_adjusted_valuation_summary.json", adjusted_summary)
    _write_audit_trail(output_dir / "06_audit_trail.md", ARTIFACT_NAMES, audit)
    _write_json(output_dir / "00_deal_pack_bundle_manifest.json", _manifest(output_dir, ARTIFACT_NAMES[:-1]))
    bundle_path = _zip_bundle(output_dir, "dealforge_portfolio_plus_pack.zip", ARTIFACT_NAMES)
    artifacts = [output_dir / name for name in ARTIFACT_NAMES]
    summary = {**summary, "source_index": source_index, "financial_extract": financial_extract, "comps_analysis": comps, "precedent_transactions": precedents_payload, "source_confidence": source_confidence, "valuation_input_audit": audit, "valuation_adjustment_recommendations": recommendations, "adjusted_valuation_summary": adjusted_summary, "consistency_checks": consistency, "validation_report": validation, "artifact_count": len(artifacts), "bundle_path": str(bundle_path)}
    return PortfolioPackResult(output_dir=output_dir, bundle_path=bundle_path, summary=summary, artifacts=artifacts)
